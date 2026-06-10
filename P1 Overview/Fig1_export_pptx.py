#!/usr/bin/env python3
"""
Fig1 互动编辑导出脚本
======================
为每个 panel 生成独立 SVG，然后组装为 PPTX。
在 PowerPoint 中打开后，每个 panel 可以独立拖拽、缩放、调整位置。
"""

from pathlib import Path
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle, FancyArrowPatch, Polygon
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import os
import tempfile

# ─── 复用主脚本的全局配置 ──────────────────────────────────────────
mpl.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans", "sans-serif"],
    "svg.fonttype": "none",
    "pdf.fonttype": 42,
    "ps.fonttype": 42,
    "font.size": 6.0,
    "axes.linewidth": 0.55,
    "axes.spines.right": False,
    "axes.spines.top": False,
    "legend.frameon": False,
    "savefig.dpi": 600,
    "figure.facecolor": "white",
})

OUT_DIR = Path(__file__).resolve().parent
SEARCH_DIRS = [Path.cwd() / "materials", OUT_DIR / "materials", Path("/mnt/data/materials"), Path.cwd(), OUT_DIR]

# 颜色常量 (与主脚本一致)
BLUE   = "#0F4D92"
BLUE2  = "#DCEBFF"
GRAY   = "#F5F5F3"
DARK   = "#1C1C1C"
TEXT   = "#222222"
MUTED  = "#6F6F6F"
GREEN  = "#2E9E44"
RED    = "#E53935"
PALE_YELLOW = "#FFF7D6"
PALE_BLUE   = "#EAF3FF"
PALE_GREEN  = "#EDF7E8"
PALE_PURPLE = "#F1ECFA"
PALE_RED    = "#FDEBEB"

# 整图尺寸 (inch)
FIG_W, FIG_H = 180/25.4, 190/25.4

# ══════════════════════════════════════════════════════════════════
# 以下 helper 函数与主脚本完全一致 (精简复制, 保证独立运行)
# ══════════════════════════════════════════════════════════════════

def find_asset(name):
    for d in SEARCH_DIRS:
        p = d / name
        if p.exists():
            return p
    return None

def safe_imread(name):
    p = find_asset(name)
    if p is None:
        return None
    try:
        return plt.imread(str(p))
    except Exception:
        return None

def hide_ax(ax):
    ax.set_xticks([]); ax.set_yticks([])
    for sp in ax.spines.values():
        sp.set_visible(False)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)

def panel_label(ax, label, x=-0.035, y=1.03):
    ax.text(x, y, label, transform=ax.transAxes, ha="left", va="top",
            fontsize=8.2, fontweight="bold", color="black")

def rounded(ax, xy, wh, fc, ec="#BBBBBB", lw=0.65, r=0.018, z=1, alpha=1.0):
    x, y = xy; w, h = wh
    patch = FancyBboxPatch((x, y), w, h,
                           boxstyle=f"round,pad=0.004,rounding_size={r}",
                           fc=fc, ec=ec, lw=lw, zorder=z, alpha=alpha)
    ax.add_patch(patch)
    return patch

def rect_ext(x, y, w, h):
    return (x, x + w, y, y + h)

def union_ext(*exts):
    exts = [e for e in exts if e is not None]
    if not exts:
        return None
    return (min(e[0] for e in exts), max(e[1] for e in exts),
            min(e[2] for e in exts), max(e[3] for e in exts))

def data_unit_display_aspect(ax):
    fig = ax.figure
    pos = ax.get_position()
    ax_w = pos.width * fig.get_figwidth()
    ax_h = pos.height * fig.get_figheight()
    x0, x1 = ax.get_xlim()
    y0, y1 = ax.get_ylim()
    return (ax_w / abs(x1 - x0)) / (ax_h / abs(y1 - y0))

def axes_size_mm(ax, width_mm, height_mm=None):
    if height_mm is None:
        height_mm = width_mm
    fig = ax.figure
    pos = ax.get_position()
    ax_w_mm = pos.width * fig.get_figwidth() * 25.4
    ax_h_mm = pos.height * fig.get_figheight() * 25.4
    x0, x1 = ax.get_xlim()
    y0, y1 = ax.get_ylim()
    return (width_mm / ax_w_mm * abs(x1 - x0),
            height_mm / ax_h_mm * abs(y1 - y0))

def centered_gap_arrow(ax, left_ext, right_ext, y, length=0.018,
                       color="#333333", lw=1.0, ms=8, zorder=8):
    if left_ext is None or right_ext is None:
        return
    gap_start, gap_end = left_ext[1], right_ext[0]
    gap = gap_end - gap_start
    if gap <= 0:
        return
    use_len = min(length, gap * 0.72)
    cx = (gap_start + gap_end) / 2
    ax.add_patch(FancyArrowPatch((cx - use_len/2, y), (cx + use_len/2, y),
                                 arrowstyle="-|>", mutation_scale=ms,
                                 lw=lw, color=color, shrinkA=0, shrinkB=0,
                                 zorder=zorder))

def magnify_connector(fig, source_ax, source_ext, target_ax, target_span, source_pad=None):
    x0, x1, y0, y1 = source_ext
    if source_pad is None:
        pad_x, pad_y = 0.009, 0.014
    else:
        pad_x, pad_y = source_pad
    source_left = (x0-pad_x, y0-pad_y)
    source_right = (x1+pad_x, y0-pad_y)
    target_left = (target_span[0], 1.006)
    target_right = (target_span[1], 1.006)
    to_fig = fig.transFigure.inverted()
    region = [
        to_fig.transform(source_ax.transData.transform(source_left)),
        to_fig.transform(source_ax.transData.transform(source_right)),
        to_fig.transform(target_ax.transAxes.transform(target_right)),
        to_fig.transform(target_ax.transAxes.transform(target_left)),
    ]
    fig.add_artist(Polygon(region, closed=True, transform=fig.transFigure,
                           fc="#E8EBED", ec="#B0B8C0", lw=0.3, alpha=1.0, zorder=0.1))
    source_ax.add_patch(Rectangle((x0-pad_x, y0-pad_y),
                                  x1-x0+2*pad_x, y1-y0+2*pad_y,
                                  fill=True, fc="#E8EBED", ec="none",
                                  alpha=1.0, zorder=2, clip_on=False))
    target_ax.patch.set_alpha(0)
    target_ax.add_patch(Rectangle((0.000, 0.000), 1.000, 1.000,
                                  transform=target_ax.transAxes,
                                  fill=True, fc="#E8EBED", ec="none",
                                  zorder=-2, clip_on=False))

def draw_image_original_ratio(ax, img, box, zorder=4, pad_frac=0.0):
    if img is None:
        return None
    x0, y0, w, h = box
    x0 += w * pad_frac
    y0 += h * pad_frac
    w *= (1 - 2 * pad_frac)
    h *= (1 - 2 * pad_frac)
    ih, iw = img.shape[:2]
    img_ar = iw / ih
    unit_ar = data_unit_display_aspect(ax)
    box_ar = (w / h) * unit_ar
    if img_ar >= box_ar:
        draw_w = w
        draw_h = w * unit_ar / img_ar
    else:
        draw_h = h
        draw_w = h * img_ar / unit_ar
    cx = x0 + w/2
    cy = y0 + h/2
    ext = (cx - draw_w/2, cx + draw_w/2, cy - draw_h/2, cy + draw_h/2)
    ax.imshow(img, extent=ext, zorder=zorder, aspect="auto")
    return ext

def draw_sensor_photo(ax, center, size, img_name, title=None, fallback_color="#444444"):
    x, y = center; w, h = size
    img = safe_imread(img_name)
    if img is not None:
        obj_ext = draw_image_original_ratio(ax, img, (x-w/2, y-h/2, w, h), zorder=4)
    else:
        rounded(ax, (x-w*0.43, y-h*0.30), (w*0.86, h*0.52), fc="#363636", ec="#151515", lw=0.7, r=0.025, z=4)
        ax.add_patch(Circle((x, y), min(w,h)*0.22, fc="#0F0F0F", ec="#666", lw=0.5, zorder=5))
        ax.add_patch(Circle((x, y), min(w,h)*0.11, fc="#202A35", ec="#999", lw=0.5, zorder=6))
        ax.add_patch(Rectangle((x-w*0.32, y-h*0.45), w*0.64, h*0.16, fc=fallback_color, ec="none", alpha=0.9, zorder=4))
        for k, col in enumerate(["#d84b4b", "#45a15f", "#4c7bd9"]*3):
            theta = np.deg2rad(35 + k*22)
            ax.add_patch(Circle((x+np.cos(theta)*w*0.33, y+np.sin(theta)*h*0.26), min(w,h)*0.018, fc=col, ec="none", zorder=7))
        obj_ext = rect_ext(x-w/2, y-h/2, w, h)
    if title:
        ax.text(x, y-h/2-0.020, title, ha="center", va="top", fontsize=5.0, color=TEXT)
    return obj_ext

def draw_raw_stack(ax, center, size):
    x, y = center; w, h = size
    rng = np.random.default_rng(7)
    img = safe_imread("raw_rgb.png")
    if img is None:
        img = rng.random((80, 80, 3)) * np.array([0.9, 0.65, 0.55])
        img[..., 1] *= rng.random((80, 80))*0.55 + 0.35
        img[..., 2] *= rng.random((80, 80))*0.7 + 0.25
    img_ext = draw_image_original_ratio(ax, img, (x-w/2, y-h/2, w, h), zorder=8)
    layer_dx, layer_dy = axes_size_mm(ax, 0.42, 0.42)
    img_w = img_ext[1] - img_ext[0]
    img_h = img_ext[3] - img_ext[2]
    layer_exts = []
    for i in range(4):
        dx = (4-i) * layer_dx
        dy = (4-i) * layer_dy
        layer_ext = rect_ext(img_ext[0] + dx, img_ext[2] + dy, img_w, img_h)
        layer_exts.append(layer_ext)
        ax.add_patch(Rectangle((layer_ext[0], layer_ext[2]), img_w, img_h,
                               fc="white", ec="#BBBBBB", lw=0.35, zorder=2+i))
    return union_ext(img_ext, *layer_exts)

def depth_array(n=96):
    yy, xx = np.mgrid[-1:1:complex(n), -1:1:complex(n)]
    z = 2.1*np.exp(-((xx*0.9)**2+(yy+0.15)**2)/0.22) - 0.25*np.exp(-((xx-0.45)**2+(yy-0.25)**2)/0.045)
    return z

def draw_depth(ax, extent, title=None, show_scale=False):
    img = safe_imread("my_depth_2d.png")
    x0, x1, y0, y1 = extent
    if img is not None:
        obj_ext = draw_image_original_ratio(ax, img, (x0, y0, x1-x0, y1-y0), zorder=3)
    else:
        ax.imshow(depth_array(), extent=extent, cmap="viridis", vmin=-2, vmax=2, zorder=3, aspect="auto")
        obj_ext = extent
    if title:
        ax.text((x0+x1)/2, y0-0.020, title, ha="center", va="top", fontsize=5.0, color=TEXT)
    if show_scale:
        grad = np.linspace(-2, 2, 100)[:, None]
        ax.imshow(grad, extent=(x1+0.012, x1+0.025, y0, y1), cmap="viridis", vmin=-2, vmax=2, zorder=3, aspect="auto")
        ax.text(x1+0.029, y1, "2", va="center", fontsize=5.0)
        ax.text(x1+0.029, (y0+y1)/2, "0", va="center", fontsize=5.0)
        ax.text(x1+0.029, y0, "-2", va="center", fontsize=5.0)
        ax.text(x1+0.020, y1+0.017, "Depth\n(mm)", ha="center", va="bottom", fontsize=5.0)
    return obj_ext

def solver_card(ax, center, size, title, stages, footer, accent, fill):
    x, y = center
    w, h = size
    rounded(ax, (x-w/2, y-h/2), (w, h), fc=fill, ec=accent, lw=0.75, r=0.012, z=3)
    ax.text(x, y+h*0.30, title, ha="center", va="center",
            fontsize=5.0, fontweight="bold", color=accent)
    stage_y = y+h*0.01
    stage_w = w * 0.205
    stage_h = h * 0.23
    stage_gap = w * 0.115
    total_w = 3*stage_w + 2*stage_gap
    start_x = x - total_w/2
    for i, stage in enumerate(stages):
        sx = start_x + i*(stage_w + stage_gap)
        rounded(ax, (sx, stage_y-stage_h/2), (stage_w, stage_h),
                fc="white", ec=accent, lw=0.38, r=0.004, z=4)
        ax.text(sx+stage_w/2, stage_y, stage, ha="center", va="center",
                fontsize=5.0, color=TEXT, zorder=5, linespacing=0.95)
        if i < 2:
            arrow_start = sx + stage_w + stage_gap*0.18
            arrow_end = sx + stage_w + stage_gap*0.82
            ax.add_patch(FancyArrowPatch((arrow_start, stage_y), (arrow_end, stage_y),
                                         arrowstyle="-|>", mutation_scale=4.8,
                                         lw=0.62, color=accent,
                                         shrinkA=0, shrinkB=0, zorder=7))
    ax.text(x, y-h*0.33, footer, ha="center", va="center",
            fontsize=5.0, color=accent, fontweight="bold")
    return rect_ext(x-w/2, y-h/2, w, h)

def metric_note(ax, center, size, heading, detail, color):
    x, y = center
    w, h = size
    ax.text(x, y+h*0.36, heading, color=color, ha="center", va="center",
            fontsize=5.0, fontweight="bold")
    if isinstance(detail, (list, tuple)):
        line_ys = [y+h*0.10, y-h*0.10, y-h*0.30]
        for item, yy in zip(detail, line_ys):
            if len(item) == 3:
                category, value, line_color = item
                line = f"{category}: {value}"
            else:
                line, line_color = item
            ax.text(x-w*0.43, yy, line, color=line_color, ha="left", va="center",
                    fontsize=5.0)
    else:
        ax.text(x, y-h*0.15, detail, color=TEXT, ha="center", va="center",
                fontsize=5.0, linespacing=1.12)
    return rect_ext(x-size[0]/2, y-h/2, size[0], h)

def leader_label(ax, anchor, label_y, text, side="right"):
    x, y = anchor
    if label_y is None:
        label_y = y
    if side == "right":
        text_x, ha = 0.640, "left"
        curve_end_x = text_x - 0.044
        line_end = text_x - 0.018
        start_x = x + 0.006
    else:
        text_x, ha = 0.120, "right"
        curve_end_x = text_x + 0.044
        line_end = text_x + 0.018
        start_x = x - 0.006
    mid_x = (start_x + curve_end_x) / 2
    ctrl_y = y * 0.58 + label_y * 0.42
    t = np.linspace(0, 1, 40)
    curve_x = (1-t)**2 * start_x + 2*(1-t)*t * mid_x + t**2 * curve_end_x
    curve_y = (1-t)**2 * y + 2*(1-t)*t * ctrl_y + t**2 * label_y
    for line_x, line_y in [
            (curve_x, curve_y),
            ([curve_end_x, line_end], [label_y, label_y])]:
        line, = ax.plot(line_x, line_y, color="#6D6D6D", lw=0.58,
                        ls=(0, (2.6, 2.2)), zorder=5,
                        solid_capstyle="round")
        line.set_dash_capstyle("round")
    ax.text(text_x, label_y, text, ha=ha, va="center", fontsize=5.0, color=TEXT)

def pipeline_box(ax, xy, wh, title, detail, fill, edge="#B7B7B7", title_color=TEXT):
    x, y = xy
    w, h = wh
    rounded(ax, (x, y), (w, h), fc=fill, ec=edge, lw=0.6, r=0.010, z=2)
    ax.text(x+w/2, y+h*0.64, title, ha="center", va="center",
            fontsize=5.0, fontweight="bold", color=title_color, linespacing=1.02)
    ax.text(x+w/2, y+h*0.27, detail, ha="center", va="center",
            fontsize=5.0, color=TEXT, linespacing=1.02)
    return rect_ext(x, y, w, h)

def image_placeholder(ax, box, title, note=None):
    x, y, w, h = box
    ax.add_patch(Rectangle((x, y), w, h, fc="#FAFBFC", ec="none", lw=0, zorder=2))
    ax.text(x+w/2, y+h/2+(0.014 if note else 0), title, ha="center", va="center",
            fontsize=5.0, color=MUTED, fontweight="bold")
    if note:
        ax.text(x+w/2, y+h/2-0.030, note, ha="center", va="center",
                fontsize=5.0, color=MUTED)
    return rect_ext(x, y, w, h)

def pixel_stream_icon(ax, box):
    x, y, w, h = box
    colors = [
        ["#E88972", "#7EA89A", "#E9CE8C", "#83A7C1"],
        ["#89B0C9", "#E99C8E", "#77A899", "#E5CF92"],
        ["#E9CC8A", "#7FA6C0", "#E88C78", "#8BB19D"],
        ["#82A5C0", "#E9C98B", "#81AC9E", "#E98976"],
    ]
    cell_w = w / 4
    cell_h = h / 4
    for row in range(4):
        for col in range(4):
            ax.add_patch(Rectangle((x + col*cell_w, y + (3-row)*cell_h),
                                   cell_w, cell_h, fc=colors[row][col],
                                   ec="#666666", lw=0.3, zorder=3))
    ax.add_patch(Rectangle((x, y), w, h, fill=False, ec="#444444", lw=0.5, zorder=4))
    return rect_ext(x, y, w, h)

def span_label(ax, start, end, y, text):
    ax.add_patch(FancyArrowPatch((start, y), (end, y), arrowstyle="<->",
                                 mutation_scale=5, lw=0.5, color="#333333",
                                 shrinkA=0, shrinkB=0))
    ax.text((start+end)/2, y-0.040, text, ha="center", va="top",
            fontsize=5.0, color=TEXT)


# ══════════════════════════════════════════════════════════════════
# 各 panel 绘制函数
# ══════════════════════════════════════════════════════════════════

def draw_panel_a(ax):
    """panel a：Host vs Near-sensor 流程对比。"""
    hide_ax(ax)
    panel_label(ax, "a", x=-0.028, y=1.015)
    ax.text(0.030, 0.945, "Host reconstruction", fontsize=6.0,
            fontweight="bold", color="#333333", ha="left", va="center")

    xs = [0.100, 0.270, 0.505, 0.705, 0.890]
    yt, yb = 0.750, 0.200
    img_size = axes_size_mm(ax, 10.0, 10.0)
    img_w, img_h = img_size

    # 上方: Host 流程
    top_exts = []
    top_exts.append(draw_sensor_photo(ax, (xs[0], yt), img_size, "GelSight.png", None))
    top_exts.append(draw_raw_stack(ax, (xs[1], yt), img_size))
    top_exts.append(solver_card(ax, (xs[2], yt), (0.215, 0.215),
                                "Iterative multigrid", ("copy", "iter.", "depth"),
                                "variable cycles", "#666666", "#FFFFFF"))
    top_exts.append(draw_depth(ax, (xs[3]-img_w/2, xs[3]+img_w/2,
                                    yt-img_h/2, yt+img_h/2), title="depth output"))
    top_exts.append(metric_note(ax, (xs[4], yt), (0.205, 0.285),
                               "variable delay",
                               [("Latency", "unstable", "#8F2E2A"),
                                ("Power", "high startup", "#B43D35"),
                                ("Energy", "low efficiency", "#D15B4E")], RED))
    for left, right in zip(top_exts[:-1], top_exts[1:]):
        centered_gap_arrow(ax, left, right, y=yt, length=0.027,
                           color="#5A5A5A", lw=0.90, ms=6.4)

    # 下方: Near-sensor 流程
    bottom_exts = []
    bottom_exts.append(draw_sensor_photo(ax, (xs[0], yb), img_size, "Sensor.png", None, fallback_color="#93C7BD"))
    bottom_exts.append(draw_raw_stack(ax, (xs[1], yb), img_size))
    bottom_exts.append(solver_card(ax, (xs[2], yb), (0.215, 0.215),
                                   "Spectral DST", ("stream", "$1/\\lambda$", "depth"),
                                   "fixed 35,107 cycles", BLUE, "#FFFFFF"))
    bottom_exts.append(draw_depth(ax, (xs[3]-img_w/2, xs[3]+img_w/2,
                                       yb-img_h/2, yb+img_h/2), title=None))
    ax.text(xs[3], yb+img_h/2+0.020, "on-chip depth", ha="center", va="bottom",
            fontsize=5.0, color=TEXT, zorder=6)
    bottom_exts.append(metric_note(ax, (xs[4], yb), (0.205, 0.285),
                                   "spectral Poisson",
                                   [("Latency", "0.211 ms, zero jitter", "#137A43"),
                                    ("Power", "0.305 W on chip", "#208F53"),
                                    ("Energy", "0.031 mJ/frame", "#35A869")], GREEN))
    bottom_union = union_ext(*bottom_exts)
    extend_down_5mm = axes_size_mm(ax, 1.0, 5.0)[1]
    bg_y0 = bottom_union[2] - 0.026 - extend_down_5mm
    bg_y1 = min(0.420, bottom_union[3] + 0.026)
    ax.add_patch(Rectangle((0.015, bg_y0), 0.970, bg_y1-bg_y0, fc="#F7FAFF",
                           ec="none", zorder=0, clip_on=False))
    ax.text(0.030, bg_y1+0.018, "Near-sensor (this work)", fontsize=5.7,
            fontweight="bold", color=BLUE, ha="left", va="center")
    for left, right in zip(bottom_exts[:-1], bottom_exts[1:]):
        centered_gap_arrow(ax, left, right, y=yb, length=0.027,
                           color=BLUE, lw=0.95, ms=6.4)
    return {"sensor": bottom_exts[0], "pipeline": bottom_exts[2]}

def draw_panel_b(ax, label="b"):
    """panel b：传感器爆炸图。"""
    hide_ax(ax)
    panel_label(ax, label, x=-0.035, y=1.02)
    img = safe_imread("Sensor_Explosion.png")
    if img is not None:
        img_ext = draw_image_original_ratio(ax, img, (0.020, 0.020, 0.504, 0.948), zorder=2)
    else:
        img_ext = image_placeholder(ax, (0.190, 0.020, 0.408, 0.948),
                                    "exploded view", "replace image")
    ix0, ix1, iy0, iy1 = img_ext
    ih = iy1 - iy0
    label_ys = np.linspace(0.860, 0.180, 7)
    part_marks = [
        (0.91, label_ys[0], "PDMS elastomer"),
        (0.80, label_ys[1], "optical window"),
        (0.68, label_ys[2], "lens holder"),
        (0.53, label_ys[3], "RGB illumination"),
        (0.39, label_ys[4], "support frame"),
        (0.25, label_ys[5], "IMX219 CMOS"),
        (0.10, label_ys[6], "enclosure"),
    ]
    for frac, label_y, text in part_marks:
        leader_label(ax, (ix1-0.012, iy0+ih*frac), label_y, text, "right")

def draw_panel_c(ax, label="c"):
    """panel c：片上 pipeline 展开图。"""
    hide_ax(ax)
    panel_label(ax, label, x=-0.028, y=1.015)
    pixel_w, pixel_h = axes_size_mm(ax, 8.0, 8.0)
    pixel_x, pixel_y = 0.020, 0.555 - pixel_h/2
    stream_ext = pixel_stream_icon(ax, (pixel_x, pixel_y, pixel_w, pixel_h))
    ax.text(pixel_x+pixel_w/2, pixel_y+pixel_h+0.025, "CMOS\npixels", ha="center", va="bottom",
            fontsize=5.0, color=TEXT, linespacing=1.05)
    ax.text(pixel_x+pixel_w/2, pixel_y-0.025, "400 FPS", ha="center", va="top",
            fontsize=5.0, color=TEXT)

    stage_specs = [
        (0.080, "Photo\nstereo", "RGB LUT\n$g_x,g_y$", PALE_YELLOW, "#C9AE62"),
        (0.072, "Div.", "$\\nabla\\cdot g$", PALE_GREEN, "#8EB27B"),
        (0.090, "DST", "row/col", PALE_PURPLE, "#A394BF"),
        (0.090, "Spectral\nsolve", "$\\hat f/\\lambda$", PALE_RED, "#D1A19A"),
        (0.090, "IDST", "inverse", PALE_PURPLE, "#A394BF"),
        (0.076, "Output", "sfix24", PALE_YELLOW, "#C9AE62"),
    ]
    depth_x = 0.895
    stage_gap = (depth_x - stream_ext[1] - sum(item[0] for item in stage_specs)) / (len(stage_specs) + 1)
    stage_x = stream_ext[1] + stage_gap
    specs = []
    for w, title, detail, fill, edge in stage_specs:
        specs.append((stage_x, w, title, detail, fill, edge))
        stage_x += w + stage_gap
    stage_y, stage_h = 0.350, 0.430
    stages = [pipeline_box(ax, (x, stage_y), (w, stage_h), title, detail, fill, edge)
              for x, w, title, detail, fill, edge in specs]
    pipe_arrow_len = 0.028
    centered_gap_arrow(ax, stream_ext, stages[0], y=0.555, length=pipe_arrow_len,
                       color="#222222", lw=0.72, ms=4.2)
    for left, right in zip(stages[:-1], stages[1:]):
        centered_gap_arrow(ax, left, right, y=0.555, length=pipe_arrow_len,
                           color="#222222", lw=0.72, ms=4.2)

    out_size = axes_size_mm(ax, 8.0, 8.0)
    depth_y = 0.555 - out_size[1]/2
    depth_ext = draw_depth(ax, (depth_x, depth_x+out_size[0], depth_y,
                                depth_y+out_size[1]), title=None)
    centered_gap_arrow(ax, stages[-1], depth_ext, y=0.555, length=pipe_arrow_len,
                       color="#222222", lw=0.72, ms=4.2)
    ax.text((depth_ext[0]+depth_ext[1])/2, depth_ext[3]+0.025, "Depth out\n256 x 256",
            ha="center", va="bottom", fontsize=5.0, color=TEXT, linespacing=1.05)
    ax.text((depth_ext[0]+depth_ext[1])/2, depth_ext[2]-0.025, "depth clock",
            ha="center", va="top", fontsize=5.0, color=TEXT)

    span_label(ax, stages[0][0], stages[1][1], 0.292, "streaming, line-buffered")
    span_label(ax, stages[2][0], stages[4][1], 0.292, "double-buffered transpose")
    span_label(ax, stages[5][0], stages[5][1], 0.292, "streaming")
    span_label(ax, stages[0][0], stages[-1][1], 0.182,
               "fixed latency: 0.211 ms")

def draw_panel_d(ax):
    """panel d：传感器在机械手上的应用图。"""
    hide_ax(ax)
    panel_label(ax, "d", x=-0.035, y=1.015)
    frames = [
        (["robot_hand_setup.png", "sensor_on_finger.png",
          "dexterous_hand_setup.png"], (0.075, 0.570, 0.850, 0.340),
         "sensor-on-hand setup"),
        (["robot_hand_grasp.png", "grasp_moment.png",
          "feather_touch.png"], (0.075, 0.085, 0.850, 0.340),
         "grasp with depth feedback"),
    ]
    for names, box, fallback in frames:
        img = None
        for name in names:
            img = safe_imread(name)
            if img is not None:
                break
        if img is None:
            image_placeholder(ax, box, fallback, "robot hand + sensor photo")
        else:
            draw_image_original_ratio(ax, img, box, zorder=3)
    ax.add_patch(FancyArrowPatch((0.500, 0.535), (0.500, 0.465),
                                 arrowstyle="-|>", mutation_scale=6,
                                 lw=0.55, color="#56616B",
                                 shrinkA=0, shrinkB=0))
    ax.text(0.525, 0.500, "0.211 ms\ndepth out", ha="left", va="center",
            fontsize=4.8, color=BLUE, fontweight="bold", linespacing=1.05)
    ax.text(0.500, 0.015, "on-chip tactile depth reflex",
            ha="center", va="bottom", fontsize=5.0, color=TEXT)

def draw_panel_e(ax):
    """panel e：雷达对比图。"""
    labels = ["Low\nlatency", "Power\nefficiency", "Throughput",
              "Board\nfootprint", "Timing\ndeterminism"]
    angles = np.linspace(0, 2*np.pi, len(labels), endpoint=False).tolist()
    loop = angles + angles[:1]
    RADAR_BLUE   = "#0F4D92"
    RADAR_CYAN   = "#008A8A"
    RADAR_ORANGE = "#D97C2B"
    RADAR_GRAY   = "#767676"
    series = [
        ("This work",       [0.97, 0.92, 0.94, 0.96, 0.98], RADAR_BLUE,   1.15),
        ("Jetson Orin NX",  [0.64, 0.52, 0.78, 0.42, 0.50], RADAR_CYAN,   0.78),
        ("GPU",             [0.74, 0.25, 0.95, 0.18, 0.43], RADAR_ORANGE, 0.78),
        ("CPU",             [0.30, 0.34, 0.36, 0.62, 0.29], RADAR_GRAY,   0.78),
    ]
    panel_label(ax, "e", x=-0.18, y=1.10)
    ax.set_theta_offset(np.pi/2)
    ax.set_theta_direction(-1)
    ax.set_ylim(0, 1.0)
    ax.set_xticks(angles)
    ax.set_xticklabels(labels, fontsize=5.0, color=TEXT)
    ax.set_yticks([0.25, 0.50, 0.75, 1.00])
    ax.set_yticklabels([])
    ax.grid(color="#D8DEE3", lw=0.42)
    ax.spines["polar"].set_color("#B6C0C8")
    ax.spines["polar"].set_linewidth(0.5)
    for name, values, color, lw in series:
        closed = values + values[:1]
        ax.plot(loop, closed, color=color, lw=lw, label=name, zorder=4)
        if "This work" in name:
            ax.fill(loop, closed, color=color, alpha=0.11, zorder=2)
    ax.legend(loc="upper center", bbox_to_anchor=(0.50, -0.14), ncol=2,
              frameon=False, fontsize=5.0, handlelength=1.5,
              columnspacing=0.9, handletextpad=0.35)
    ax.text(0.50, -0.28, "normalised per dimension to best observed value\nabsolute values in Extended Data Table",
            transform=ax.transAxes, ha="center", va="top", fontsize=4.5, color=MUTED)

def draw_panel_f(ax):
    """panel f：按压时间序列 + 深度 colorbar。"""
    hide_ax(ax)
    panel_label(ax, "f", x=-0.035, y=1.015)
    times = ["0.000", "0.133", "0.267", "0.400", "0.533", "0.667"]
    stages = ["baseline", "first contact", "indent", "spreading", "peak load", "hold"]
    descriptions = ["no load", "touch onset", "local dent",
                    "contact grows", "maximum indent", "steady hold"]
    files = [f"depth_sequence/depth_{i:03d}.png" for i in range(6)]
    press_files = [f"press_images/press_image_{i:02d}.png" for i in range(6)]
    w, h = axes_size_mm(ax, 12.0, 12.0)
    gap = (0.920 - 6*w) / 5
    positions = [(0.040 + i*(w+gap), 0.520) for i in range(6)]
    progresses = [0.00, 0.12, 0.38, 0.58, 0.78, 0.95]
    for t, stage, desc, name, press_name, (x, y), progress in zip(
            times, stages, descriptions, files, press_files, positions, progresses):
        img = safe_imread(name)
        if img is None:
            image_placeholder(ax, (x, y, w, h), "depth map", "measured frame")
        else:
            draw_image_original_ratio(ax, img, (x, y, w, h), zorder=3)
        ax.text(x+w/2, y+h+0.035, f"{t} s", ha="center", va="bottom", fontsize=5.0)
        ax.add_patch(FancyArrowPatch((x+w/2, y-0.018), (x+w/2, 0.205+h+0.020),
                                     arrowstyle="-|>", mutation_scale=4.2,
                                     lw=0.42, color="#78848D",
                                     shrinkA=0, shrinkB=0))
        ax.text(x+w/2, 0.452, stage, ha="center", va="center",
                fontsize=5.0, color=TEXT)
        press_img = safe_imread(press_name)
        if press_img is None:
            image_placeholder(ax, (x, 0.205, w, h), "press image", "replace")
        else:
            draw_image_original_ratio(ax, press_img, (x, 0.205, w, h), zorder=3)
        ax.text(x+w/2, 0.105, desc, ha="center", va="center",
                fontsize=5.0, color=MUTED)
    for i in range(len(positions) - 1):
        x0 = positions[i][0] + w + 0.008
        x1 = positions[i+1][0] - 0.008
        for yy in (positions[i][1] + h/2, 0.205 + h/2):
            if x1 > x0:
                ax.add_patch(FancyArrowPatch((x0, yy), (x1, yy),
                                             arrowstyle="-|>", mutation_scale=4.2,
                                             lw=0.42, color="#8A949D",
                                             shrinkA=0, shrinkB=0))
    # 全光谱深度标尺: 蓝(0) → 绿→黄→橙 → 红(2 mm)
    from matplotlib.colors import LinearSegmentedColormap, Normalize
    from matplotlib.colorbar import ColorbarBase
    depth_cmap = LinearSegmentedColormap.from_list("depth_mm", [
        "#0B1F5C", "#1A4B8C", "#2E7DB5", "#3FA37E", "#6EBF4C",
        "#D5C93C", "#F0A328", "#E05A1F", "#B2182B",
    ], N=256)
    depth_norm = Normalize(vmin=0, vmax=2)
    cbar_ax = ax.inset_axes([0.035, -0.085, 0.930, 0.022])
    cb = ColorbarBase(cbar_ax, cmap=depth_cmap, norm=depth_norm, orientation="horizontal",
                      ticks=[0, 0.5, 1.0, 1.5, 2.0])
    cb.ax.tick_params(labelsize=4.8, colors=TEXT, width=0.4)
    cb.outline.set_linewidth(0.4)
    cbar_ax.text(1.015, 0.5, "Depth (mm)", transform=cbar_ax.transAxes,
                 ha="left", va="center", fontsize=4.8, color=TEXT)


# ══════════════════════════════════════════════════════════════════
# 独立 panel 渲染 + SVG 导出
# ══════════════════════════════════════════════════════════════════

PANEL_SPECS = {
    "a": (draw_panel_a, (8.0, 2.0)),          # 宽 panel
    "b": (draw_panel_b, (3.0, 3.0)),          # 窄高 inset
    "c": (draw_panel_c, (6.5, 3.0)),          # 宽 pipeline
    "d": (draw_panel_d, (3.0, 4.0)),          # 竖高应用图
    "e": (draw_panel_e, (4.0, 4.0)),          # 雷达图 (polar)
    "f": (draw_panel_f, (9.0, 3.5)),          # 宽时间序列
}

def render_panel(panel_key, out_dir):
    """渲染单个 panel 为独立 SVG + 高分辨率 PNG 文件。"""
    draw_fn, (w_in, h_in) = PANEL_SPECS[panel_key]
    if panel_key == "e":
        fig = plt.figure(figsize=(w_in, h_in))
        ax = fig.add_subplot(111, projection="polar")
    else:
        fig, ax = plt.subplots(figsize=(w_in, h_in))

    draw_fn(ax)

    # SVG — 矢量源文件 (Illustrator/Inkscape 可编辑)
    svg_path = out_dir / f"panel_{panel_key}.svg"
    fig.savefig(svg_path, bbox_inches="tight", pad_inches=0.02, transparent=False,
                facecolor="white", edgecolor="none")
    # PNG — 高分辨率光栅 (用于 PPTX 插入)
    png_path = out_dir / f"panel_{panel_key}.png"
    fig.savefig(png_path, dpi=600, bbox_inches="tight", pad_inches=0.02,
                facecolor="white", edgecolor="none")
    plt.close(fig)
    return svg_path, png_path


# ══════════════════════════════════════════════════════════════════
# PPTX 组装
# ══════════════════════════════════════════════════════════════════

# PPTX 画布 (mm) — A4 横向，留边距
SLIDE_W_MM = 297.0
SLIDE_H_MM = 210.0
MARGIN_MM = 8.0

# 每个 panel 在 PPTX 中的位置和大小 (left_mm, top_mm, width_mm)
# 基于原图 make_figure 的 add_axes 布局换算
# 原图总尺寸: 180mm x 190mm
FIG_W_MM = 180.0
FIG_H_MM = 195.0

PANEL_LAYOUT_PPTX = {
    # key: (left_mm, top_mm, width_mm, height_mm)
    "a": (MARGIN_MM + 0.040 * FIG_W_MM,    MARGIN_MM + 0.800 * FIG_H_MM,    0.650 * FIG_W_MM, 0.160 * FIG_H_MM),
    "b": (MARGIN_MM + 0.040 * FIG_W_MM,    MARGIN_MM + 0.575 * FIG_H_MM,    0.205 * FIG_W_MM, 0.190 * FIG_H_MM),
    "c": (MARGIN_MM + 0.265 * FIG_W_MM,    MARGIN_MM + 0.575 * FIG_H_MM,    0.425 * FIG_W_MM, 0.190 * FIG_H_MM),
    "d": (MARGIN_MM + 0.730 * FIG_W_MM,    MARGIN_MM + 0.615 * FIG_H_MM,    0.230 * FIG_W_MM, 0.345 * FIG_H_MM),
    "e": (MARGIN_MM + 0.060 * FIG_W_MM,    MARGIN_MM + 0.298 * FIG_H_MM,    0.220 * FIG_W_MM, 0.257 * FIG_H_MM),
    "f": (MARGIN_MM + 0.335 * FIG_W_MM,    MARGIN_MM + 0.301 * FIG_H_MM,    0.625 * FIG_W_MM, 0.239 * FIG_H_MM),
}

def mm_to_emu(mm):
    return int(mm * 36000)  # 1 mm = 36000 EMU

def build_pptx(svg_dir, out_path):
    """将独立 panel SVG 组装为可编辑 PPTX。"""
    prs = Presentation()
    prs.slide_width  = mm_to_emu(SLIDE_W_MM)
    prs.slide_height = mm_to_emu(SLIDE_H_MM)

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # 白色背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    panel_order = ["a", "b", "c", "d", "e", "f"]
    for key in panel_order:
        png_file = svg_dir / f"panel_{key}.png"
        if not png_file.exists():
            print(f"  WARNING: {png_file} not found, skipping panel {key}")
            continue

        left_mm, top_mm, w_mm, h_mm = PANEL_LAYOUT_PPTX[key]

        left = mm_to_emu(left_mm)
        top  = mm_to_emu(top_mm)
        width  = mm_to_emu(w_mm)
        height = mm_to_emu(h_mm)

        pic = slide.shapes.add_picture(
            str(png_file),
            left, top, width, height
        )
        pic.name = f"panel_{key}"

        print(f"  panel_{key}: ({left_mm:.0f}, {top_mm:.0f}) mm  {w_mm:.0f}x{h_mm:.0f} mm")

    prs.save(str(out_path))
    print(f"\nPPTX saved: {out_path}")


# ══════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    svg_dir = OUT_DIR / "panel_svgs"
    svg_dir.mkdir(exist_ok=True)

    print("=== Rendering individual panels (SVG + PNG) ===")
    for key in ["a", "b", "c", "d", "e", "f"]:
        svg_path, png_path = render_panel(key, svg_dir)
        print(f"  panel_{key}: {svg_path.name}  +  {png_path.name}")

    print("\n=== Building interactive PPTX ===")
    pptx_path = OUT_DIR / "Fig1_Nature_Sensors_Interactive.pptx"
    build_pptx(svg_dir, pptx_path)
    print(f"\nDone: {pptx_path}")
    print("Usage:")
    print("  - PowerPoint: 每个 panel 是独立对象，可拖动/缩放/替换")
    print("  - Illustrator: 打开 panel_svgs/*.svg 进行元素级矢量编辑")
    print("  - Reference:   原始完整图 Fig1_nature_sensors_artguide.{svg,pdf}")
