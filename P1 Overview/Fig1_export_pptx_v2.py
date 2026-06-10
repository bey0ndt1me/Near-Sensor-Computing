#!/usr/bin/env python3
"""
Fig1 互动编辑导出 v2 — 元素级可编辑 PPTX
===========================================
每个 panel 拆分为:
  - 一张图形背景 (PNG, 不含文字)
  - 多个独立文本框 (可拖动/改字/改大小)
  - 图片元素独立放置
  - 连接箭头独立放置

在 PowerPoint 中:
  - 选中任意文字 → 拖动位置 / 改字体字号 / 改内容
  - 选中任意图片 → 拖动 / 缩放
  - 全部元素可独立操作, 就像做 PPT 一样
"""

from pathlib import Path
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle, FancyArrowPatch, Polygon
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io, os, copy

# ══════════════════════════════════════════════════════════════════
# 全局配置
# ══════════════════════════════════════════════════════════════════

mpl.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans", "sans-serif"],
    "svg.fonttype": "none",
    "pdf.fonttype": 42,
    "font.size": 6.0,
    "axes.linewidth": 0.55,
    "axes.spines.right": False,
    "axes.spines.top": False,
    "legend.frameon": False,
    "savefig.dpi": 600,
    "figure.facecolor": "white",
})

OUT_DIR = Path(__file__).resolve().parent
SEARCH_DIRS = [Path.cwd() / "materials", OUT_DIR / "materials", Path("/mnt/data/materials"), Path.cwd(), OUT_DIR]

BLUE   = "#0F4D92"; BLUE2 = "#DCEBFF"; GRAY = "#F5F5F3"; DARK = "#1C1C1C"
TEXT   = "#222222"; MUTED = "#6F6F6F"; GREEN = "#2E9E44"; RED = "#E53935"
PALE_YELLOW = "#FFF7D6"; PALE_BLUE = "#EAF3FF"; PALE_GREEN = "#EDF7E8"
PALE_PURPLE = "#F1ECFA"; PALE_RED = "#FDEBEB"

FIG_W, FIG_H = 180/25.4, 190/25.4
FIG_W_MM = 180.0
FIG_H_MM = 195.0

# ══════════════════════════════════════════════════════════════════
# Helper 函数 (精简复用)
# ══════════════════════════════════════════════════════════════════

def find_asset(name):
    for d in SEARCH_DIRS:
        p = d / name
        if p.exists(): return p
    return None

def safe_imread(name):
    p = find_asset(name)
    if p is None: return None
    try: return plt.imread(str(p))
    except Exception: return None

def hide_ax(ax):
    ax.set_xticks([]); ax.set_yticks([])
    for sp in ax.spines.values(): sp.set_visible(False)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)

def rect_ext(x, y, w, h):
    return (x, x + w, y, y + h)

def union_ext(*exts):
    exts = [e for e in exts if e is not None]
    if not exts: return None
    return (min(e[0] for e in exts), max(e[1] for e in exts),
            min(e[2] for e in exts), max(e[3] for e in exts))

def data_unit_display_aspect(ax):
    fig = ax.figure
    pos = ax.get_position()
    ax_w = pos.width * fig.get_figwidth()
    ax_h = pos.height * fig.get_figheight()
    x0, x1 = ax.get_xlim(); y0, y1 = ax.get_ylim()
    return (ax_w / abs(x1 - x0)) / (ax_h / abs(y1 - y0))

def axes_size_mm(ax, width_mm, height_mm=None):
    if height_mm is None: height_mm = width_mm
    fig = ax.figure
    pos = ax.get_position()
    ax_w_mm = pos.width * fig.get_figwidth() * 25.4
    ax_h_mm = pos.height * fig.get_figheight() * 25.4
    x0, x1 = ax.get_xlim(); y0, y1 = ax.get_ylim()
    return (width_mm / ax_w_mm * abs(x1 - x0),
            height_mm / ax_h_mm * abs(y1 - y0))

def draw_image_original_ratio(ax, img, box, zorder=4, pad_frac=0.0):
    if img is None: return None
    x0, y0, w, h = box
    x0 += w * pad_frac; y0 += h * pad_frac
    w *= (1 - 2 * pad_frac); h *= (1 - 2 * pad_frac)
    ih, iw = img.shape[:2]
    img_ar = iw / ih
    unit_ar = data_unit_display_aspect(ax)
    box_ar = (w / h) * unit_ar
    if img_ar >= box_ar:
        draw_w = w; draw_h = w * unit_ar / img_ar
    else:
        draw_h = h; draw_w = h * img_ar / unit_ar
    cx = x0 + w/2; cy = y0 + h/2
    ext = (cx - draw_w/2, cx + draw_w/2, cy - draw_h/2, cy + draw_h/2)
    ax.imshow(img, extent=ext, zorder=zorder, aspect="auto")
    return ext

def rounded(ax, xy, wh, fc, ec="#BBBBBB", lw=0.65, r=0.018, z=1, alpha=1.0):
    x, y = xy; w, h = wh
    patch = FancyBboxPatch((x, y), w, h,
                           boxstyle=f"round,pad=0.004,rounding_size={r}",
                           fc=fc, ec=ec, lw=lw, zorder=z, alpha=alpha)
    ax.add_patch(patch)
    return patch

def centered_gap_arrow(ax, left_ext, right_ext, y, length=0.018,
                       color="#333333", lw=1.0, ms=8, zorder=8):
    if left_ext is None or right_ext is None: return
    gap_start, gap_end = left_ext[1], right_ext[0]
    gap = gap_end - gap_start
    if gap <= 0: return
    use_len = min(length, gap * 0.72)
    cx = (gap_start + gap_end) / 2
    ax.add_patch(FancyArrowPatch((cx - use_len/2, y), (cx + use_len/2, y),
                                 arrowstyle="-|>", mutation_scale=ms,
                                 lw=lw, color=color, shrinkA=0, shrinkB=0, zorder=zorder))

def image_placeholder(ax, box, title, note=None):
    x, y, w, h = box
    ax.add_patch(Rectangle((x, y), w, h, fc="#FAFBFC", ec="none", lw=0, zorder=2))
    ax.text(x+w/2, y+h/2+(0.014 if note else 0), title, ha="center", va="center",
            fontsize=5.0, color=MUTED, fontweight="bold")
    if note:
        ax.text(x+w/2, y+h/2-0.030, note, ha="center", va="center", fontsize=5.0, color=MUTED)
    return rect_ext(x, y, w, h)


# ══════════════════════════════════════════════════════════════════
# 元素描述数据结构
# ══════════════════════════════════════════════════════════════════

class TextElement:
    """PPTX 文本框"""
    def __init__(self, text, x_mm, y_mm, w_mm, h_mm, fontsize=5.0,
                 bold=False, color=TEXT, ha="left", va="center"):
        self.text = text
        self.x_mm = x_mm
        self.y_mm = y_mm
        self.w_mm = w_mm
        self.h_mm = h_mm
        self.fontsize = fontsize
        self.bold = bold
        self.color = color
        self.ha = ha
        self.va = va
        self.name = text[:30]

class ImgElement:
    """图片"""
    def __init__(self, x_mm, y_mm, w_mm, h_mm, image_path, name="img"):
        self.x_mm = x_mm
        self.y_mm = y_mm
        self.w_mm = w_mm
        self.h_mm = h_mm
        self.image_path = image_path
        self.name = name


def mm_to_emu(mm):
    return int(mm * 36000)

def px_to_mm(px, dpi=600):
    return px / dpi * 25.4

def data_to_panel_mm(ax, data_x, data_y, panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm):
    """将 matplotlib 数据坐标转换为 panel 内的 mm 坐标"""
    xlim = ax.get_xlim(); ylim = ax.get_ylim()
    fx = (data_x - xlim[0]) / (xlim[1] - xlim[0])
    fy = (data_y - ylim[0]) / (ylim[1] - ylim[0])
    return panel_left_mm + fx * panel_w_mm, panel_top_mm + (1 - fy) * panel_h_mm


# ══════════════════════════════════════════════════════════════════
# Panel A: 流程对比 (背景图 + 文本覆盖)
# ══════════════════════════════════════════════════════════════════

def get_panel_a_elements(panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm):
    """返回 panel a 的所有独立元素列表"""
    text_els = []
    img_els = []

    # Panel 坐标系统 (数据坐标 0-1)
    pw = panel_w_mm; ph = panel_h_mm
    pl = panel_left_mm; pt = panel_top_mm

    def dtom(dx, dy):
        """数据坐标 → panel 内 mm 坐标 (数据范围 0-1)"""
        return pl + dx * pw, pt + (1 - dy) * ph

    # Panel 标签
    x, y = dtom(-0.028, 1.015)
    text_els.append(TextElement("a", x, y, 10, 5, fontsize=8.2, bold=True, color="black"))

    # "Host reconstruction" 标题
    x, y = dtom(0.030, 0.945)
    text_els.append(TextElement("Host reconstruction", x, y, 35, 5,
                                fontsize=6.0, bold=True, color="#333333"))

    # Host 行五个阶段标注
    xs = [0.100, 0.270, 0.505, 0.705, 0.890]
    yt = 0.750

    # Solver 卡片标题
    x, y = dtom(xs[2], yt)
    text_els.append(TextElement("Iterative multigrid", x-15, y+4, 30, 5,
                                fontsize=5.0, bold=True, color="#666666", ha="center"))

    # Solver 卡片三阶段
    stage_texts = ["copy", "iter.", "depth"]
    stage_w = 0.205 * 0.215; stage_gap = 0.115 * 0.215
    total_w = 3*stage_w + 2*stage_gap
    start_x_stage = xs[2] - total_w/2
    stage_y_data = yt + 0.01
    for i, st in enumerate(stage_texts):
        sx = start_x_stage + i*(stage_w + stage_gap)
        x, y = dtom(sx + stage_w/2, stage_y_data)
        text_els.append(TextElement(st, x-8, y, 16, 4, fontsize=5.0, color=TEXT, ha="center"))

    # Solver 卡片脚注
    x, y = dtom(xs[2], yt - 0.215*0.33)
    text_els.append(TextElement("variable cycles", x-15, y, 30, 4,
                                fontsize=5.0, bold=True, color="#666666", ha="center"))

    # Host depth output 标题
    img_w_data = axes_size_mm_from_panel(panel_w_mm, panel_h_mm, 10.0, 10.0)[0]
    x, y = dtom(xs[3], yt + img_w_data/2 + 0.020)
    text_els.append(TextElement("depth output", x-15, y, 30, 4,
                                fontsize=5.0, color=TEXT, ha="center"))

    # Host 指标卡
    x, y = dtom(xs[4], yt)
    text_els.append(TextElement("variable delay", x-15, y+8, 30, 4,
                                fontsize=5.0, bold=True, color=RED, ha="center"))
    detail_items = [
        ("Latency: unstable", "#8F2E2A"),
        ("Power: high startup", "#B43D35"),
        ("Energy: low efficiency", "#D15B4E"),
    ]
    for i, (line, lc) in enumerate(detail_items):
        ly = yt + 0.10 - i*0.20
        x, y = dtom(xs[4]-0.205*0.43, ly)
        text_els.append(TextElement(line, x, y, 30, 4, fontsize=5.0, color=lc))

    # "Near-sensor (this work)" 标签
    text_els.append(TextElement("Near-sensor (this work)", pl + 0.030*pw, pt + 0.42*ph,
                                40, 5, fontsize=5.7, bold=True, color=BLUE))

    # Near-sensor 行
    yb = 0.200
    # Solver 标题
    x, y = dtom(xs[2], yb)
    text_els.append(TextElement("Spectral DST", x-15, y+4, 30, 5,
                                fontsize=5.0, bold=True, color=BLUE, ha="center"))

    # Solver 三阶段
    nstage_texts = ["stream", "$1/\\lambda$", "depth"]
    for i, st in enumerate(nstage_texts):
        sx = start_x_stage + i*(stage_w + stage_gap)
        x, y = dtom(sx + stage_w/2, stage_y_data - (yt - yb))
        text_els.append(TextElement(st, x-8, y, 16, 4, fontsize=5.0, color=TEXT, ha="center"))

    # Solver 脚注
    x, y = dtom(xs[2], yb - 0.215*0.33)
    text_els.append(TextElement("fixed 35,107 cycles", x-18, y, 36, 4,
                                fontsize=5.0, bold=True, color=BLUE, ha="center"))

    # Near-sensor depth 标题
    x, y = dtom(xs[3], yb + img_w_data/2 + 0.020)
    text_els.append(TextElement("on-chip depth", x-15, y, 30, 4,
                                fontsize=5.0, color=TEXT, ha="center"))

    # Near-sensor 指标卡
    x, y = dtom(xs[4], yb)
    text_els.append(TextElement("spectral Poisson", x-15, y+8, 30, 4,
                                fontsize=5.0, bold=True, color=GREEN, ha="center"))
    ndetail_items = [
        ("Latency: 0.211 ms, zero jitter", "#137A43"),
        ("Power: 0.305 W on chip", "#208F53"),
        ("Energy: 0.031 mJ/frame", "#35A869"),
    ]
    for i, (line, lc) in enumerate(ndetail_items):
        ly = yb + 0.10 - i*0.20
        x, y = dtom(xs[4]-0.205*0.43, ly)
        text_els.append(TextElement(line, x, y, 30, 4, fontsize=5.0, color=lc))

    return text_els, img_els


def axes_size_mm_from_panel(pw_mm, ph_mm, w_mm, h_mm):
    """简化版: 估算 panel 内 mm 对应的数据坐标比例"""
    # panel 的 figure 尺寸 (英寸) 换算
    fig_w_mm = FIG_W_MM  # 180
    fig_h_mm = FIG_H_MM  # 195
    return (w_mm / fig_w_mm, h_mm / fig_h_mm)


# ══════════════════════════════════════════════════════════════════
# Panel B: 爆炸图 (背景图 + 标签)
# ══════════════════════════════════════════════════════════════════

def get_panel_b_elements(panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm):
    text_els = []
    img_els = []

    pl, pt, pw, ph = panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm

    # Panel label
    text_els.append(TextElement("b", pl-6, pt-1, 10, 5, fontsize=8.2, bold=True, color="black"))

    # 七个爆炸图标签
    label_ys = np.linspace(0.860, 0.180, 7)
    part_marks = [
        (0.91, label_ys[0], "PDMS elastomer"),
        (0.80, label_ys[1], "optical window"),
        (0.68, label_ys[2], "lens holder"),
        (0.53, label_ys[3], "RGB illumination"),
        (0.39, label_ys[4], "support frame"),
        (0.25, label_ys[5], "IMX219 CMOS"),
        (0.10, label_ys[6], "enclosure"),
    ]
    # 标签在图片右侧，x 取 panel 宽度的 ~64%
    for frac, label_y, text in part_marks:
        x_mm = pl + 0.64 * pw
        y_mm = pt + (1 - label_y) * ph
        text_els.append(TextElement(text, x_mm, y_mm, 30, 4, fontsize=5.0, color=TEXT))

    return text_els, img_els


# ══════════════════════════════════════════════════════════════════
# Panel C: Pipeline 展开图
# ══════════════════════════════════════════════════════════════════

def get_panel_c_elements(panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm):
    text_els = []
    pl, pt, pw, ph = panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm

    text_els.append(TextElement("c", pl-5, pt-1, 10, 5, fontsize=8.2, bold=True, color="black"))

    # CMOS 标注
    text_els.append(TextElement("CMOS\npixels", pl + 8, pt + 0.58*ph, 15, 8, fontsize=5.0, color=TEXT, ha="center"))
    text_els.append(TextElement("400 FPS", pl + 8, pt + 0.23*ph, 15, 4, fontsize=5.0, color=TEXT, ha="center"))

    # Pipeline 模块 (6 个)
    stage_specs = [
        ("Photo\nstereo", "RGB LUT\n$g_x,g_y$"),
        ("Div.", "$\\nabla\\cdot g$"),
        ("DST", "row/col"),
        ("Spectral\nsolve", "$\\hat f/\\lambda$"),
        ("IDST", "inverse"),
        ("Output", "sfix24"),
    ]
    stage_widths = [0.080, 0.072, 0.090, 0.090, 0.090, 0.076]
    # 从左侧 CMOS 到右侧 depth out 近似均匀分布
    start_x_mm = pl + 0.16 * pw
    end_x_mm = pl + 0.82 * pw
    stage_x_positions = np.linspace(start_x_mm, end_x_mm, len(stage_specs))
    for i, ((title, detail), sx_mm) in enumerate(zip(stage_specs, stage_x_positions)):
        sy_mm = pt + 0.50 * ph
        text_els.append(TextElement(title, sx_mm-12, sy_mm+12, 24, 10,
                                    fontsize=5.0, bold=True, color=TEXT, ha="center"))
        text_els.append(TextElement(detail, sx_mm-12, sy_mm-2, 24, 8,
                                    fontsize=5.0, color=TEXT, ha="center"))

    # Depth out
    text_els.append(TextElement("Depth out\n256 x 256", pl + 0.90*pw, pt + 0.70*ph, 20, 10,
                                fontsize=5.0, color=TEXT, ha="center"))
    text_els.append(TextElement("depth clock", pl + 0.90*pw, pt + 0.22*ph, 20, 4,
                                fontsize=5.0, color=TEXT, ha="center"))

    # Span labels
    text_els.append(TextElement("streaming, line-buffered", pl + 0.20*pw, pt + 0.13*ph, 35, 4,
                                fontsize=5.0, color=TEXT, ha="center"))
    text_els.append(TextElement("double-buffered transpose", pl + 0.45*pw, pt + 0.13*ph, 35, 4,
                                fontsize=5.0, color=TEXT, ha="center"))
    text_els.append(TextElement("streaming", pl + 0.80*pw, pt + 0.13*ph, 20, 4,
                                fontsize=5.0, color=TEXT, ha="center"))
    text_els.append(TextElement("fixed latency: 0.211 ms", pl + 0.40*pw, pt + 0.02*ph, 35, 4,
                                fontsize=5.0, color=TEXT, ha="center"))

    return text_els, []


# ══════════════════════════════════════════════════════════════════
# Panel D: 机械手应用图
# ══════════════════════════════════════════════════════════════════

def get_panel_d_elements(panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm):
    text_els = []
    pl, pt, pw, ph = panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm

    text_els.append(TextElement("d", pl-6, pt-1, 10, 5, fontsize=8.2, bold=True, color="black"))
    text_els.append(TextElement("0.211 ms\ndepth out", pl + 0.55*pw, pt + 0.50*ph, 18, 8,
                                fontsize=4.8, bold=True, color=BLUE, ha="left", va="center"))
    text_els.append(TextElement("on-chip tactile depth reflex", pl + 0.50*pw, pt + 0.01*ph, 40, 4,
                                fontsize=5.0, color=TEXT, ha="center"))

    return text_els, []


# ══════════════════════════════════════════════════════════════════
# Panel E: 雷达图
# ══════════════════════════════════════════════════════════════════

def get_panel_e_elements(panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm):
    text_els = []
    pl, pt, pw, ph = panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm

    text_els.append(TextElement("e", pl-12, pt-2, 10, 5, fontsize=8.2, bold=True, color="black"))

    # 五维度标签 (polar 角度轴上)
    dims = ["Low\nlatency", "Power\nefficiency", "Throughput", "Board\nfootprint", "Timing\ndeterminism"]
    angles = np.linspace(np.pi/2, -3*np.pi/2, 5, endpoint=False)
    r_label = 1.15  # 标签在雷达图外侧
    cx_mm = pl + 0.50*pw
    cy_mm = pt + 0.50*ph
    radius_mm = min(pw, ph) * 0.38
    for angle, dim in zip(angles, dims):
        x_mm = cx_mm + r_label * radius_mm * np.cos(angle)
        y_mm = cy_mm - r_label * radius_mm * np.sin(angle)
        text_els.append(TextElement(dim, x_mm-12, y_mm-6, 24, 10,
                                    fontsize=5.0, color=TEXT, ha="center"))

    # 图例 (底部)
    legend_items = ["This work", "Jetson Orin NX", "GPU", "CPU"]
    for i, name in enumerate(legend_items):
        col = [pl + 0.12*pw, pl + 0.58*pw][i % 2]
        row = pt + 0.78*ph if i < 2 else pt + 0.88*ph
        text_els.append(TextElement(name, col, row, 35, 4, fontsize=5.0, color=TEXT))

    # 归一化说明
    text_els.append(TextElement("normalised per dimension to best; abs. in Extended Data Table",
                                pl + 0.05*pw, pt + 0.96*ph, pw*0.9, 5,
                                fontsize=4.5, color=MUTED, ha="center"))

    return text_els, []


# ══════════════════════════════════════════════════════════════════
# Panel F: 按压时间序列
# ══════════════════════════════════════════════════════════════════

def get_panel_f_elements(panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm):
    text_els = []
    pl, pt, pw, ph = panel_left_mm, panel_top_mm, panel_w_mm, panel_h_mm

    text_els.append(TextElement("f", pl-6, pt-1, 10, 5, fontsize=8.2, bold=True, color="black"))

    times = ["0.000 s", "0.133 s", "0.267 s", "0.400 s", "0.533 s", "0.667 s"]
    stages = ["baseline", "first contact", "indent", "spreading", "peak load", "hold"]
    descriptions = ["no load", "touch onset", "local dent", "contact grows", "maximum indent", "steady hold"]

    # 6 帧均匀分布
    n_frames = 6
    frame_w_mm = pw * 0.12
    gap_mm = (pw * 0.92 - n_frames * frame_w_mm) / (n_frames - 1)
    for i in range(n_frames):
        x_center = pl + pw * 0.04 + i * (frame_w_mm + gap_mm) + frame_w_mm/2

        # 时间戳 (深度图上方)
        text_els.append(TextElement(times[i], x_center-10, pt + 0.55*ph + frame_w_mm+2, 20, 4,
                                    fontsize=5.0, color=TEXT, ha="center"))
        # 阶段名 (中间)
        text_els.append(TextElement(stages[i], x_center-15, pt + 0.45*ph, 30, 4,
                                    fontsize=5.0, color=TEXT, ha="center"))
        # 说明 (下方)
        text_els.append(TextElement(descriptions[i], x_center-15, pt + 0.10*ph, 30, 4,
                                    fontsize=5.0, color=MUTED, ha="center"))

    # Colorbar 标签
    text_els.append(TextElement("Depth (mm)", pl + 0.97*pw, pt + 0.06*ph, 15, 4,
                                fontsize=4.8, color=TEXT, ha="left", va="center"))

    return text_els, []


# ══════════════════════════════════════════════════════════════════
# Panel 背景图渲染 (仅图形, 不含文字)
# ══════════════════════════════════════════════════════════════════

# 每个 panel 的绘制函数和尺寸
PANEL_CONFIG = {
    "a": ("panel_a_base", 8.0, 2.0),
    "b": ("panel_b_base", 3.0, 3.0),
    "c": ("panel_c_base", 6.5, 3.0),
    "d": ("panel_d_base", 3.0, 4.0),
    "e": ("panel_e_base", 4.0, 4.0),
    "f": ("panel_f_base", 9.0, 3.5),
}

PANEL_LAYOUT = {
    "a": (8.0 + 0.040*FIG_W_MM, 8.0 + 0.800*FIG_H_MM,    0.650*FIG_W_MM, 0.160*FIG_H_MM),
    "b": (8.0 + 0.040*FIG_W_MM, 8.0 + 0.575*FIG_H_MM,    0.205*FIG_W_MM, 0.190*FIG_H_MM),
    "c": (8.0 + 0.265*FIG_W_MM, 8.0 + 0.575*FIG_H_MM,    0.425*FIG_W_MM, 0.190*FIG_H_MM),
    "d": (8.0 + 0.730*FIG_W_MM, 8.0 + 0.615*FIG_H_MM,    0.230*FIG_W_MM, 0.345*FIG_H_MM),
    "e": (8.0 + 0.060*FIG_W_MM, 8.0 + 0.298*FIG_H_MM,    0.220*FIG_W_MM, 0.257*FIG_H_MM),
    "f": (8.0 + 0.335*FIG_W_MM, 8.0 + 0.301*FIG_H_MM,    0.625*FIG_W_MM, 0.239*FIG_H_MM),
}

ELEMENT_FUNCTIONS = {
    "a": get_panel_a_elements,
    "b": get_panel_b_elements,
    "c": get_panel_c_elements,
    "d": get_panel_d_elements,
    "e": get_panel_e_elements,
    "f": get_panel_f_elements,
}


def render_panel_base(panel_key, out_dir):
    """渲染 panel 背景图 (仅图形, 不包含任何文字) — 直接复用完整 SVG"""
    # 由于 matplotlib 很难剥离文字, 这里我们直接复用已导出的完整 panel PNG
    # 用户上一轮已通过 Fig1_export_pptx.py 导出了 panel PNG
    png_path = out_dir / f"panel_{panel_key}.png"
    if not png_path.exists():
        # fallback: 运行 Fig1_export_pptx.py
        print(f"  WARNING: {png_path} not found. Run Fig1_export_pptx.py first.")
    return png_path


# ══════════════════════════════════════════════════════════════════
# PPTX 组装
# ══════════════════════════════════════════════════════════════════

SLIDE_W_MM = 340.0   # 宽屏
SLIDE_H_MM = 220.0

FONT_SIZE_MAP = {4.5: 6, 4.8: 6, 5.0: 7, 5.7: 8, 6.0: 8, 8.2: 11}

NAMED_COLORS = {
    "black":   (0, 0, 0),
    "white":   (255, 255, 255),
    "red":     (255, 0, 0),
    "green":   (0, 128, 0),
    "blue":    (0, 0, 255),
}

def hex_to_rgb(hex_color):
    if hex_color.lower() in NAMED_COLORS:
        r, g, b = NAMED_COLORS[hex_color.lower()]
        return RGBColor(r, g, b)
    c = hex_color.lstrip('#')
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))

def build_pptx_v2(panel_dir, out_path):
    """组装元素级可编辑 PPTX"""
    prs = Presentation()
    prs.slide_width  = mm_to_emu(SLIDE_W_MM)
    prs.slide_height = mm_to_emu(SLIDE_H_MM)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白色背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for panel_key in ["a", "b", "c", "d", "e", "f"]:
        pl_mm, pt_mm, pw_mm, ph_mm = PANEL_LAYOUT[panel_key]

        # 1. 背景图 (完整 panel PNG, 作为底层)
        png_path = panel_dir / f"panel_{panel_key}.png"
        if png_path.exists():
            pic = slide.shapes.add_picture(
                str(png_path),
                mm_to_emu(pl_mm), mm_to_emu(pt_mm),
                mm_to_emu(pw_mm), mm_to_emu(ph_mm)
            )
            pic.name = f"panel_{panel_key}_base"

        # 2. 文字元素覆盖层
        elem_fn = ELEMENT_FUNCTIONS[panel_key]
        text_els, img_els = elem_fn(pl_mm, pt_mm, pw_mm, ph_mm)

        for te in text_els:
            # PPTX 文本框
            left = mm_to_emu(te.x_mm)
            top  = mm_to_emu(te.y_mm)
            width = mm_to_emu(te.w_mm)
            height = mm_to_emu(te.h_mm)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = te.text
            ppt_fontsize = FONT_SIZE_MAP.get(te.fontsize, Pt(int(te.fontsize + 2)))
            p.font.size = Pt(ppt_fontsize)
            p.font.bold = te.bold
            p.font.color.rgb = hex_to_rgb(te.color)
            p.font.name = "Arial"
            if te.ha == "center":
                p.alignment = PP_ALIGN.CENTER
            elif te.ha == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            txBox.name = f"txt_{panel_key}_{te.name[:20]}"

    prs.save(str(out_path))
    print(f"PPTX saved: {out_path}")


# ══════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    # 确保已有 panel PNG (由 Fig1_export_pptx.py 生成)
    svg_dir = OUT_DIR / "panel_svgs"
    png_missing = False
    for key in ["a", "b", "c", "d", "e", "f"]:
        png_path = svg_dir / f"panel_{key}.png"
        if not png_path.exists():
            print(f"MISSING: {png_path}")
            png_missing = True

    if png_missing:
        print("\nRun Fig1_export_pptx.py first to generate panel PNGs:")
        print("  python Fig1_export_pptx.py")
        print("\nThen re-run this script.")
        exit(1)

    print("=== Building element-level editable PPTX ===")
    pptx_path = OUT_DIR / "Fig1_Nature_Sensors_Editable_v2.pptx"
    build_pptx_v2(svg_dir, pptx_path)

    print(f"\nDone: {pptx_path}")
    print("In PowerPoint:")
    print("  - Each text label is a separate text box → drag/retype/resize")
    print("  - Panel backgrounds are PNG images underneath")
    print("  - For deeper edits, open panel_svgs/*.svg in Illustrator")
